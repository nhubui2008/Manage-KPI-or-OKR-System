#!/usr/bin/env python3
"""Build the local NEXTGEN graduation presentation.

Requires python-pptx and Pillow. The generated deck is intentionally based on
the verified local-code/SQL snapshot documented on 2026-08-11; it does not make
production-deployment claims.
"""

from __future__ import annotations

import os
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
DOCS = ROOT / "docs"
OUT = DOCS / "Presentation_BaoCao_TotNghiep_NEXTGEN.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Noto Sans"

NAVY = "0B1220"
NAVY_2 = "111C33"
INK = "0F172A"
SLATE = "475569"
MUTED = "64748B"
CANVAS = "F6F7FB"
WHITE = "FFFFFF"
BORDER = "DDE3EE"
INDIGO = "4F46E5"
VIOLET = "7C3AED"
CYAN = "06B6D4"
TEAL = "0F766E"
GREEN = "16A34A"
AMBER = "F59E0B"
ROSE = "E11D48"
BLUE_LIGHT = "EEF2FF"
CYAN_LIGHT = "ECFEFF"
GREEN_LIGHT = "ECFDF5"
AMBER_LIGHT = "FFFBEB"
ROSE_LIGHT = "FFF1F2"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def add_rect(slide, x, y, w, h, fill, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    shape.line.width = Pt(0.8)
    return shape


def add_line(slide, x1, y1, x2, y2, color=BORDER, width=1.2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0,
    font=FONT,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_multiline(slide, lines, x, y, w, h, size=16, color=INK, gap=6):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.02)
    frame.margin_top = frame.margin_bottom = Inches(0.02)
    for i, line in enumerate(lines):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
    return box


def add_bullet(slide, text, x, y, w, color=INK, accent=INDIGO, size=16):
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y + 0.08), Inches(0.12), Inches(0.12)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = rgb(accent)
    dot.line.color.rgb = rgb(accent)
    add_text(slide, text, x + 0.24, y, w - 0.24, 0.52, size=size, color=color)


def add_pill(slide, text, x, y, w, fill=BLUE_LIGHT, color=INDIGO, size=10):
    pill = add_rect(slide, x, y, w, 0.34, fill, fill, True)
    pill.line.width = Pt(0)
    add_text(
        slide,
        text.upper(),
        x,
        y + 0.01,
        w,
        0.3,
        size=size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def set_bg(slide, color=CANVAS):
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = rgb(color)


def add_slide_header(slide, number, eyebrow, title, subtitle=None, dark=False):
    ink = WHITE if dark else INK
    muted = "CBD5E1" if dark else MUTED
    add_pill(
        slide,
        eyebrow,
        0.65,
        0.42,
        min(3.4, max(1.55, len(eyebrow) * 0.085)),
        fill=NAVY_2 if dark else BLUE_LIGHT,
        color=CYAN if dark else INDIGO,
        size=9,
    )
    add_text(slide, f"{number:02d}", 12.18, 0.37, 0.5, 0.32, 11, muted, True, PP_ALIGN.RIGHT)
    add_text(slide, title, 0.65, 0.92, 12.0, 0.58, 27, ink, True)
    if subtitle:
        add_text(slide, subtitle, 0.65, 1.48, 11.8, 0.42, 12.5, muted)


def add_footer(slide, text="NEXTGEN · BIZEN AI-NATIVE"):
    add_line(slide, 0.65, 7.12, 12.68, 7.12, BORDER, 0.7)
    add_text(slide, text, 0.65, 7.18, 5.0, 0.18, 8.5, MUTED, True)


def add_picture_crop(slide, path, x, y, w, h):
    path = str(path)
    with Image.open(path) as image:
        iw, ih = image.size
    picture_ratio = iw / ih
    box_ratio = w / h
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    if picture_ratio > box_ratio:
        visible = box_ratio / picture_ratio
        crop = (1 - visible) / 2
        pic.crop_left = crop
        pic.crop_right = crop
    elif picture_ratio < box_ratio:
        visible = picture_ratio / box_ratio
        crop = (1 - visible) / 2
        pic.crop_top = crop
        pic.crop_bottom = crop
    return pic


def add_picture_frame(slide, path, x, y, w, h, border=BORDER):
    add_rect(slide, x - 0.04, y - 0.04, w + 0.08, h + 0.08, WHITE, border, True)
    return add_picture_crop(slide, path, x, y, w, h)


def add_card(slide, title, body, x, y, w, h, accent=INDIGO, fill=WHITE, number=None):
    card = add_rect(slide, x, y, w, h, fill, BORDER, True)
    card.line.width = Pt(0.8)
    add_rect(slide, x, y, 0.07, h, accent, accent, False).line.width = Pt(0)
    if number:
        add_pill(slide, number, x + 0.22, y + 0.2, 0.48, fill=BLUE_LIGHT, color=accent, size=9)
        tx = x + 0.82
        tw = w - 1.04
    else:
        tx = x + 0.28
        tw = w - 0.52
    add_text(slide, title, tx, y + 0.2, tw, 0.34, 14, INK, True)
    add_text(slide, body, x + 0.28, y + 0.67, w - 0.52, h - 0.82, 11.5, SLATE)
    return card


def add_metric(slide, value, label, x, y, w, accent=INDIGO, note=None):
    add_rect(slide, x, y, w, 1.15, WHITE, BORDER, True)
    add_text(slide, value, x + 0.22, y + 0.16, w - 0.44, 0.42, 24, accent, True)
    add_text(slide, label, x + 0.22, y + 0.62, w - 0.44, 0.24, 11, INK, True)
    if note:
        add_text(slide, note, x + 0.22, y + 0.88, w - 0.44, 0.18, 8.5, MUTED)


def new_slide(prs, bg=CANVAS):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, bg)
    return slide


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    # 01 · Cover
    slide = new_slide(prs, NAVY)
    add_rect(slide, 8.9, -0.6, 5.2, 5.2, VIOLET, VIOLET, True).rotation = 18
    add_rect(slide, 10.4, 4.6, 3.5, 3.5, CYAN, CYAN, True).rotation = 18
    add_text(slide, "BÁO CÁO DỰ ÁN TỐT NGHIỆP", 0.78, 0.64, 5.8, 0.38, 13, CYAN, True)
    add_text(
        slide,
        "HỆ THỐNG HỖ TRỢ\nVẬN HÀNH THÔNG MINH",
        0.78,
        1.3,
        7.8,
        1.65,
        31,
        WHITE,
        True,
    )
    add_text(
        slide,
        "Quản lý đa cấp · KPI/OKR liên thông · Hỗ trợ quyết định bằng AI",
        0.82,
        3.12,
        7.2,
        0.62,
        17,
        "CBD5E1",
    )
    add_rect(slide, 0.82, 4.1, 6.5, 1.2, NAVY_2, "26334D", True)
    add_text(slide, "NHÓM NEXTGEN", 1.1, 4.34, 2.5, 0.34, 17, WHITE, True)
    add_text(slide, "GVHD: Phan Hoàng Khải", 1.1, 4.78, 3.3, 0.28, 11.5, "CBD5E1")
    add_pill(slide, "ASP.NET CORE MVC", 4.18, 4.31, 1.55, NAVY, CYAN, 8.5)
    add_pill(slide, "AI-NATIVE", 5.82, 4.31, 1.16, NAVY, CYAN, 8.5)
    add_text(slide, "Hà Nội · 2026", 0.82, 6.78, 2.0, 0.3, 10.5, "94A3B8")
    add_picture_crop(slide, DOCS / "img" / "logo.jpg", 9.02, 1.18, 3.55, 2.42)

    # 02 · Team
    slide = new_slide(prs)
    add_slide_header(slide, 2, "ĐỘI NGŨ", "Nhóm NEXTGEN — 7 thành viên, một sản phẩm liên thông")
    members = [
        ("Phạm Trần Anh Quân", "TB01758", "Leader · AI", "quan.jpg"),
        ("Phạm Trần An An", "TB01817", "Backend KPI/OKR", "an.jpg"),
        ("Bùi Nguyễn Anh Như", "TB01785", "Frontend KPI/OKR", "nhu.jpg"),
        ("Trần Thanh Phong", "TB01649", "Nền tảng · Hệ thống", "phong.png"),
        ("Nguyễn Thế Bảo", "TB01573", "Vận hành · Kanban", "bao.png"),
        ("Đoàn Quốc Khánh", "TB01544", "QA · Kiểm thử", "khanh.jpg"),
        ("Vũ Hoàng Huy Nhật", "TB01605", "Frontend vận hành", "nhat.png"),
    ]
    positions = [(0.7 + i * 3.12, 2.08) for i in range(4)] + [(2.26 + i * 3.12, 4.55) for i in range(3)]
    for (name, code, role, image), (x, y) in zip(members, positions):
        add_rect(slide, x, y, 2.75, 1.98, WHITE, BORDER, True)
        add_picture_crop(slide, DOCS / "img" / image, x + 0.18, y + 0.2, 0.78, 0.78)
        add_text(slide, name, x + 1.08, y + 0.18, 1.48, 0.5, 12.2, INK, True)
        add_text(slide, code, x + 1.08, y + 0.7, 1.4, 0.22, 9, MUTED, True)
        add_pill(slide, role, x + 0.18, y + 1.23, 2.38, BLUE_LIGHT, INDIGO, 8.5)
    add_footer(slide, "NEXTGEN · TỔ CHỨC THEO PHÂN HỆ, TÍCH HỢP TRÊN CÙNG KIẾN TRÚC")

    # 03 · Problem
    slide = new_slide(prs)
    add_slide_header(
        slide,
        3,
        "BÀI TOÁN",
        "SME cần dữ liệu xuyên suốt từ chiến lược đến thực thi",
        "Quy trình rời rạc khiến quản lý phản ứng chậm và khó kiểm chứng quyết định.",
    )
    pains = [
        ("Dữ liệu phân tán", "KPI/OKR nằm ở nhiều bảng tính và phòng ban; tổng hợp chậm, dễ lệch số.", ROSE, ROSE_LIGHT),
        ("Đánh giá cảm tính", "Thiếu baseline đo lường và workflow duyệt rõ trách nhiệm giữa các cấp.", AMBER, AMBER_LIGHT),
        ("Chiến lược đứt gãy", "Mục tiêu công ty không liên kết rõ tới KR, KPI và đầu việc hằng ngày.", INDIGO, BLUE_LIGHT),
        ("Can thiệp quá muộn", "Không có luồng check-in, cảnh báo và bằng chứng kịp thời để hỗ trợ quyết định.", TEAL, GREEN_LIGHT),
    ]
    for i, (title, body, accent, fill) in enumerate(pains):
        x = 0.75 + (i % 2) * 6.15
        y = 2.24 + (i // 2) * 2.0
        add_card(slide, title, body, x, y, 5.72, 1.62, accent, fill, str(i + 1))
    add_rect(slide, 0.75, 6.3, 11.88, 0.56, NAVY, NAVY, True)
    add_text(
        slide,
        "Mấu chốt: một nguồn dữ liệu thống nhất, quyền truy cập đúng phạm vi và AI chỉ đóng vai trò cố vấn.",
        1.0,
        6.43,
        11.4,
        0.28,
        13,
        WHITE,
        True,
        PP_ALIGN.CENTER,
    )
    add_footer(slide)

    # 04 · Solution value
    slide = new_slide(prs)
    add_slide_header(slide, 4, "GIẢI PHÁP", "Bizen kết nối quản trị mục tiêu thành một vòng vận hành khép kín")
    steps = [
        ("01", "Chiến lược", "Sứ mệnh · Mục tiêu"),
        ("02", "OKR / KR", "Phân rã đa cấp"),
        ("03", "KPI", "Target · Trọng số"),
        ("04", "Thực thi", "Task · Check-in"),
        ("05", "Đánh giá", "Duyệt · Xếp hạng"),
    ]
    for i, (n, title, body) in enumerate(steps):
        x = 0.73 + i * 2.5
        accent = [INDIGO, VIOLET, CYAN, TEAL, GREEN][i]
        add_rect(slide, x, 2.18, 2.12, 1.55, WHITE, BORDER, True)
        add_pill(slide, n, x + 0.18, 2.39, 0.42, BLUE_LIGHT, accent, 9)
        add_text(slide, title, x + 0.18, 2.86, 1.75, 0.3, 14, INK, True)
        add_text(slide, body, x + 0.18, 3.24, 1.75, 0.25, 10, MUTED)
        if i < 4:
            add_text(slide, "→", x + 2.15, 2.72, 0.32, 0.35, 20, "94A3B8", True, PP_ALIGN.CENTER)
    add_card(slide, "Một dữ liệu — nhiều góc nhìn", "Dashboard theo kỳ, phạm vi vai trò và tiến độ thực tế; giảm tổng hợp thủ công.", 0.75, 4.36, 3.72, 1.45, INDIGO, BLUE_LIGHT)
    add_card(slide, "Workflow có trách nhiệm", "Nhân viên báo cáo, quản lý duyệt, hệ thống tính toán; mọi bước có trạng thái rõ ràng.", 4.8, 4.36, 3.72, 1.45, TEAL, GREEN_LIGHT)
    add_card(slide, "AI gắn vào quyết định", "Cố vấn dùng nguồn được phép; con người sửa và xác nhận qua form nghiệp vụ chuẩn.", 8.85, 4.36, 3.72, 1.45, VIOLET, "F5F3FF")
    add_footer(slide)

    # 05 · Roles and workflow
    slide = new_slide(prs)
    add_slide_header(slide, 5, "NGƯỜI DÙNG", "Một hệ thống, năm vai trò và phạm vi dữ liệu rõ ràng")
    roles = [
        ("Admin", "Cấu hình tổ chức & quyền"),
        ("Director", "Chiến lược & phê duyệt"),
        ("Manager", "Giao KPI & review"),
        ("Employee", "Thực thi & check-in"),
        ("HR", "Kỳ đánh giá & kết quả"),
    ]
    for i, (role, body) in enumerate(roles):
        x = 0.77 + i * 2.5
        add_rect(slide, x, 2.0, 2.12, 1.08, WHITE, BORDER, True)
        add_text(slide, role, x + 0.18, 2.2, 1.76, 0.28, 13.5, INDIGO, True)
        add_text(slide, body, x + 0.18, 2.58, 1.76, 0.22, 9.7, SLATE)
    add_line(slide, 1.0, 3.7, 12.2, 3.7, "C7D2FE", 3.0)
    events = [
        (1.15, "Thiết lập", "Mục tiêu / kỳ"),
        (3.35, "Giao", "OKR / KPI / task"),
        (5.55, "Thực thi", "Kanban / check-in"),
        (7.75, "Kiểm soát", "Review / phê duyệt"),
        (9.95, "Ra quyết định", "Dashboard / đánh giá"),
    ]
    for idx, (x, title, body) in enumerate(events):
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(3.49), Inches(0.42), Inches(0.42))
        c.fill.solid(); c.fill.fore_color.rgb = rgb(INDIGO); c.line.color.rgb = rgb(WHITE); c.line.width = Pt(2)
        add_text(slide, str(idx + 1), x, 3.52, 0.42, 0.28, 9, WHITE, True, PP_ALIGN.CENTER)
        add_text(slide, title, x - 0.25, 4.08, 1.45, 0.28, 13, INK, True, PP_ALIGN.CENTER)
        add_text(slide, body, x - 0.36, 4.47, 1.68, 0.38, 10.2, MUTED, False, PP_ALIGN.CENTER)
    add_rect(slide, 0.78, 5.45, 11.82, 0.88, NAVY, NAVY, True)
    add_text(slide, "RBAC + phạm vi tenant/phòng ban/nhân viên đi cùng mọi thao tác đọc và ghi", 1.05, 5.72, 11.25, 0.32, 14, WHITE, True, PP_ALIGN.CENTER)
    add_footer(slide)

    # 06 · Architecture
    slide = new_slide(prs, NAVY)
    add_slide_header(slide, 6, "KIẾN TRÚC", "Nền tảng web phân tầng, bảo mật dữ liệu theo tenant", "ASP.NET Core MVC · EF Core · SQL Server · AI model gateway", dark=True)
    layers = [
        ("TRẢI NGHIỆM", "Razor + Bootstrap · Dashboard · OKR/KPI · Kanban · Check-in", "18243D", CYAN),
        ("NGHIỆP VỤ", "Controllers · Services · Validators · Authorization filters", "17233B", INDIGO),
        ("DỮ LIỆU", "EF Core · Transactions · RowVersion · SQL Server RLS", "152239", TEAL),
        ("AI-NATIVE", "Advisor contracts · AgentRun · Outbox · Citation metadata · RAG ACL", "132137", VIOLET),
    ]
    for i, (name, body, fill, accent) in enumerate(layers):
        y = 2.2 + i * 1.05
        add_rect(slide, 1.0, y, 11.32, 0.78, fill, "2D3A55", True)
        add_pill(slide, name, 1.25, y + 0.2, 1.45, NAVY, accent, 8.5)
        add_text(slide, body, 2.98, y + 0.2, 8.85, 0.32, 13, WHITE, True)
    add_text(slide, "Mọi write chính thức đều đi qua validator + quyền + transaction của nghiệp vụ", 1.0, 6.52, 11.3, 0.34, 13.5, "CBD5E1", True, PP_ALIGN.CENTER)

    # 07 · Dashboard
    slide = new_slide(prs)
    add_slide_header(slide, 7, "SẢN PHẨM", "Dashboard hợp nhất tiến độ, con người và hành động")
    add_picture_frame(slide, DOCS / "screenshot_dashboard.png", 0.72, 1.9, 11.9, 4.63)
    add_pill(slide, "THEO KỲ", 1.02, 2.16, 0.92, NAVY, WHITE, 8)
    add_pill(slide, "OKR / KPI", 9.15, 2.16, 1.06, NAVY, WHITE, 8)
    add_pill(slide, "AI ADVISORY", 10.36, 2.16, 1.4, NAVY, CYAN, 8)
    add_footer(slide, "GIAO DIỆN THỰC TẾ · DASHBOARD THEO PHẠM VI NGƯỜI DÙNG")

    # 08 · OKR/KPI
    slide = new_slide(prs)
    add_slide_header(slide, 8, "MỤC TIÊU", "Từ Objective đến KPI có target, trọng số và người chịu trách nhiệm")
    add_picture_frame(slide, DOCS / "screenshot_okrs.png", 0.68, 1.92, 5.92, 2.91)
    add_picture_frame(slide, DOCS / "screenshot_kpis.png", 6.74, 1.92, 5.92, 2.91)
    add_pill(slide, "OKR / KEY RESULTS", 0.9, 2.12, 1.55, NAVY, WHITE, 8)
    add_pill(slide, "KPI / ASSIGNMENT", 6.96, 2.12, 1.55, NAVY, WHITE, 8)
    flow = ["Objective", "Key Result", "KPI", "Assignee", "Check-in"]
    for i, label in enumerate(flow):
        x = 1.0 + i * 2.42
        add_rect(slide, x, 5.28, 1.72, 0.66, BLUE_LIGHT if i < 3 else GREEN_LIGHT, BORDER, True)
        add_text(slide, label, x, 5.49, 1.72, 0.24, 11.5, INDIGO if i < 3 else TEAL, True, PP_ALIGN.CENTER)
        if i < len(flow) - 1:
            add_text(slide, "→", x + 1.77, 5.47, 0.56, 0.28, 18, "94A3B8", True, PP_ALIGN.CENTER)
    add_text(slide, "Liên kết được lưu trong mô hình dữ liệu và kiểm tra lại ở server — không phụ thuộc dữ liệu browser.", 0.82, 6.32, 11.7, 0.34, 12.5, SLATE, True, PP_ALIGN.CENTER)
    add_footer(slide)

    # 09 · Check-in
    slide = new_slide(prs)
    add_slide_header(slide, 9, "THỰC THI", "Check-in có lịch sử, công thức server và cổng phê duyệt của con người")
    add_picture_frame(slide, DOCS / "screenshot_checkin.png", 0.7, 1.96, 7.72, 4.6)
    items = [
        ("1", "Nhân viên báo cáo", "Giá trị đạt được + rào cản"),
        ("2", "Server tính tiến độ", "Target, trọng số, KPI thuận/nghịch"),
        ("3", "Quản lý review", "Điểm, nhận xét, approve/reject"),
        ("4", "Snapshot chính thức", "Chỉ Approved mới làm baseline"),
    ]
    for i, (n, title, body) in enumerate(items):
        y = 2.03 + i * 1.07
        add_card(slide, title, body, 8.72, y, 3.86, 0.91, INDIGO if i < 3 else GREEN, WHITE, n)
    add_rect(slide, 8.72, 6.22, 3.86, 0.48, GREEN_LIGHT, GREEN_LIGHT, True)
    add_text(slide, "Human approval là cổng cuối", 8.72, 6.35, 3.86, 0.22, 11, GREEN, True, PP_ALIGN.CENTER)
    add_footer(slide)

    # 10 · AI principle
    slide = new_slide(prs, NAVY)
    add_slide_header(slide, 10, "AI-NATIVE", "AI đề xuất, con người quyết định", "Nguyên tắc xuyên suốt mọi advisor và workflow", dark=True)
    add_picture_frame(slide, DOCS / "screenshot_ai.png", 0.72, 2.02, 6.55, 3.73, border="334155")
    stages = [
        ("AI", "Sinh draft có cấu trúc\nvà source ID", VIOLET),
        ("SERVER", "Kiểm schema · quyền\nnguồn · fingerprint", CYAN),
        ("HUMAN", "Xem · sửa · chấp nhận\nhoặc từ chối", GREEN),
        ("WORKFLOW", "Form chuẩn mới ghi\ndữ liệu chính thức", INDIGO),
    ]
    for i, (tag, body, accent) in enumerate(stages):
        y = 2.05 + i * 1.12
        add_rect(slide, 7.68, y, 4.86, 0.88, NAVY_2, "334155", True)
        add_pill(slide, tag, 7.9, y + 0.24, 0.82, NAVY, accent, 8.2)
        add_text(slide, body, 8.94, y + 0.13, 3.28, 0.57, 12.2, WHITE, True)
    add_text(slide, "Không AI nào tự duyệt KPI/OKR, tự sửa điểm, xếp hạng hay thưởng.", 0.82, 6.3, 11.7, 0.34, 13.5, "CBD5E1", True, PP_ALIGN.CENTER)

    # 11 · Nine flows
    slide = new_slide(prs)
    add_slide_header(slide, 11, "PHẠM VI AI", "Chín luồng cố vấn đã hội tụ về cùng một kiến trúc kiểm soát")
    advisors = [
        ("Goal Planning", "Ba task plan có nguồn"),
        ("Check-in Evaluator", "Định lượng + rubric"),
        ("KR Advisor", "Đánh giá candidate KR"),
        ("Review Draft", "Nháp nhận xét 1-on-1"),
        ("Customer Segment", "Phân khúc tham khảo"),
        ("Performance", "Insight từ check-in duyệt"),
        ("KPI Suggestion", "3–5 KPI draft hợp lệ"),
        ("Chat Advisor", "Hỏi đáp có citation"),
        ("KR Suggestion", "Gợi ý/refine KR"),
    ]
    for i, (title, body) in enumerate(advisors):
        row, col = divmod(i, 3)
        x = 0.75 + col * 4.08
        y = 2.02 + row * 1.38
        accent = [INDIGO, VIOLET, CYAN, TEAL, GREEN, AMBER, INDIGO, VIOLET, CYAN][i]
        add_rect(slide, x, y, 3.72, 1.08, WHITE, BORDER, True)
        add_pill(slide, f"{i+1:02d}", x + 0.18, y + 0.18, 0.45, BLUE_LIGHT, accent, 8)
        add_text(slide, title, x + 0.78, y + 0.16, 2.68, 0.3, 12.4, INK, True)
        add_text(slide, body, x + 0.78, y + 0.57, 2.68, 0.24, 9.8, MUTED)
    add_rect(slide, 0.75, 6.34, 11.88, 0.42, BLUE_LIGHT, BLUE_LIGHT, True)
    add_text(slide, "Mẫu chung: authorized snapshot → strict JSON → validate → recheck source → human action", 0.92, 6.45, 11.5, 0.22, 11.4, INDIGO, True, PP_ALIGN.CENTER)
    add_footer(slide)

    # 12 · Evidence/RAG
    slide = new_slide(prs)
    add_slide_header(slide, 12, "BẰNG CHỨNG", "Advisor chỉ dùng nguồn được cấp quyền và có thể truy vết")
    add_picture_frame(slide, DOCS / "sequence_ai.png", 6.22, 1.92, 6.42, 4.67)
    points = [
        ("Authorized snapshot", "Server dựng scope kỳ, OKR/KPI và tenant."),
        ("Strict contract", "Model chỉ trả JSON + source IDs được cấp."),
        ("Citation metadata", "Lưu ID/hash để audit, không lưu prompt thô."),
        ("Stale-safe", "Dựng lại snapshot trong transaction trước commit."),
        ("Abstain", "Thiếu bằng chứng thì trả rỗng, không đoán."),
    ]
    for i, (title, body) in enumerate(points):
        y = 2.0 + i * 0.88
        add_text(slide, f"{i+1:02d}", 0.78, y + 0.03, 0.45, 0.24, 10, INDIGO, True)
        add_text(slide, title, 1.38, y, 2.0, 0.26, 12.3, INK, True)
        add_text(slide, body, 1.38, y + 0.32, 4.42, 0.32, 9.8, SLATE)
        add_line(slide, 0.78, y + 0.75, 5.76, y + 0.75, BORDER, 0.7)
    add_footer(slide, "RAG PIPELINE: PRIVATE BLOB → MINERU → BGE-M3 → AZURE AI SEARCH → ACL FILTER")

    # 13 · Reliability & tenant
    slide = new_slide(prs, NAVY)
    add_slide_header(slide, 13, "ĐỘ TIN CẬY", "Bảo vệ dữ liệu và xử lý bền vững ngay trong kiến trúc", dark=True)
    add_metric(slide, "57", "bảng có SQL Server RLS", 0.82, 2.0, 2.45, CYAN, "filter + block predicate")
    add_metric(slide, "100%", "scope do server xác định", 3.5, 2.0, 2.45, VIOLET, "không nhận TenantId từ browser")
    add_metric(slide, "1", "human gate trước write", 6.18, 2.0, 2.45, GREEN, "form nghiệp vụ chuẩn")
    add_metric(slide, "0", "prompt/raw response lưu mới", 8.86, 2.0, 3.05, AMBER, "chỉ metadata kiểm toán")
    left = [
        "Tenant query filters + composite tenant foreign keys",
        "SESSION_CONTEXT fail-closed cho raw SQL / IgnoreQueryFilters",
        "RAG ACL theo user · role · department",
    ]
    right = [
        "Durable outbox: lease · retry · dead-letter",
        "Idempotency + source fingerprint + row-version",
        "Serializable recheck cho dữ liệu và quyền thay đổi",
    ]
    add_rect(slide, 0.82, 3.55, 5.65, 2.52, NAVY_2, "334155", True)
    add_pill(slide, "TENANT DEFENSE", 1.08, 3.82, 1.52, NAVY, CYAN, 8.5)
    for i, item in enumerate(left):
        add_bullet(slide, item, 1.08, 4.33 + i * 0.54, 5.0, WHITE, CYAN, 11.5)
    add_rect(slide, 6.72, 3.55, 5.65, 2.52, NAVY_2, "334155", True)
    add_pill(slide, "RUNTIME SAFETY", 6.98, 3.82, 1.52, NAVY, VIOLET, 8.5)
    for i, item in enumerate(right):
        add_bullet(slide, item, 6.98, 4.33 + i * 0.54, 5.0, WHITE, VIOLET, 11.5)
    add_text(slide, "Security và reliability là điều kiện vận hành, không phải lớp trang trí sau cùng.", 0.82, 6.52, 11.55, 0.32, 13, "CBD5E1", True, PP_ALIGN.CENTER)

    # 14 · Data model
    slide = new_slide(prs)
    add_slide_header(slide, 14, "DỮ LIỆU", "Mô hình quan hệ giữ liên kết từ tổ chức đến kết quả")
    add_picture_frame(slide, DOCS / "erd_diagram.png", 0.72, 1.92, 9.2, 4.84)
    add_card(slide, "Tổ chức", "SystemUser · Employee · Department · Position", 10.18, 1.96, 2.4, 1.16, TEAL, GREEN_LIGHT)
    add_card(slide, "Mục tiêu", "OKR · Key Result · KPI · KPI Detail", 10.18, 3.28, 2.4, 1.16, INDIGO, BLUE_LIGHT)
    add_card(slide, "Thực thi", "WorkProject · WorkItem · Check-in", 10.18, 4.6, 2.4, 1.16, VIOLET, "F5F3FF")
    add_rect(slide, 10.18, 5.95, 2.4, 0.73, NAVY, NAVY, True)
    add_text(slide, "Dữ liệu liên thông\n→ báo cáo nhất quán", 10.18, 6.1, 2.4, 0.38, 11, WHITE, True, PP_ALIGN.CENTER)
    add_footer(slide)

    # 15 · Verification
    slide = new_slide(prs)
    add_slide_header(slide, 15, "KIỂM CHỨNG", "Build, test và SQL Server thật: bằng chứng chất lượng", "Snapshot local được ghi nhận ngày 11/08/2026")
    add_metric(slide, "520/520", "test solution chạy xanh", 0.78, 2.08, 2.75, GREEN, "unit + integration + SQL tests")
    add_metric(slide, "0", "warning khi build", 3.79, 2.08, 2.75, INDIGO, "solution build")
    add_metric(slide, "0", "error khi build", 6.8, 2.08, 2.75, INDIGO, "solution build")
    add_metric(slide, "57", "bảng RLS được kiểm tra", 9.81, 2.08, 2.75, CYAN, "real SQL Server")
    checks = [
        ("Migration lifecycle", "Database rỗng → latest → down → reapply; snapshot không drift."),
        ("Tenant isolation", "Raw SQL, cross-tenant insert/update và pooled connection đều bị kiểm soát."),
        ("Concurrency", "Proposal, rubric, draft và double-confirm hội tụ theo transaction/idempotency."),
        ("AI contracts", "Strict schema, citation, abstain, stale-source và human-decision được regression test."),
    ]
    for i, (title, body) in enumerate(checks):
        x = 0.78 + (i % 2) * 6.0
        y = 3.62 + (i // 2) * 1.35
        add_card(slide, title, body, x, y, 5.62, 1.08, GREEN if i == 0 else INDIGO, WHITE, "✓")
    add_rect(slide, 0.78, 6.44, 11.82, 0.38, GREEN_LIGHT, GREEN_LIGHT, True)
    add_text(slide, "Bằng chứng xác nhận phạm vi mã nguồn và SQL Server qua kiểm thử tự động.", 0.98, 6.53, 11.42, 0.22, 10.3, TEAL, True, PP_ALIGN.CENTER)
    add_footer(slide)

    # 16 · Closing
    slide = new_slide(prs, NAVY)
    add_pill(slide, "KẾT LUẬN", 0.78, 0.66, 1.2, NAVY_2, CYAN, 9)
    add_text(slide, "Một nền tảng quản trị hiệu suất\ncó thể giải thích và kiểm soát", 0.78, 1.22, 8.3, 1.28, 29, WHITE, True)
    pillars = [
        ("01", "Liên thông", "Chiến lược → OKR/KR → KPI → task/check-in → đánh giá", INDIGO),
        ("02", "Có bằng chứng", "Advisor dùng nguồn được phép, citation và source fingerprint", CYAN),
        ("03", "Con người quyết định", "AI chỉ tạo draft; workflow chuẩn giữ quyền phê duyệt", GREEN),
    ]
    for i, (num, title, body, accent) in enumerate(pillars):
        x = 0.8 + i * 4.1
        add_rect(slide, x, 3.25, 3.68, 1.72, NAVY_2, "334155", True)
        add_pill(slide, num, x + 0.24, 3.49, 0.48, NAVY, accent, 8.5)
        add_text(slide, title, x + 0.86, 3.45, 2.4, 0.32, 15, WHITE, True)
        add_text(slide, body, x + 0.24, 4.05, 3.18, 0.58, 11.2, "CBD5E1")
    add_text(slide, "XIN CẢM ƠN HỘI ĐỒNG", 0.78, 5.72, 7.2, 0.46, 22, WHITE, True)
    add_text(slide, "Nhóm NEXTGEN sẵn sàng trình bày demo và trả lời phản biện.", 0.8, 6.24, 7.6, 0.32, 13, "CBD5E1")
    add_picture_crop(slide, DOCS / "img" / "logo.jpg", 9.42, 5.25, 2.7, 1.48)

    return prs


def assert_deck(prs: Presentation) -> None:
    assert len(prs.slides) == 16, f"Expected 16 slides, got {len(prs.slides)}"
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                all_text.append(shape.text.lower())
    joined = "\n".join(all_text)
    assert "520/520" in joined, "Verification metric is missing"


def main() -> int:
    prs = build_deck()
    assert_deck(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"created={OUT}")
    print(f"slides={len(prs.slides)}")
    print(f"bytes={OUT.stat().st_size}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
